# -*- coding: utf-8 -*-
"""Pre-submission check for FathomWave-Deck.pptx.

Run this before exporting to PDF, and again on 15 September:

    python preflight.py

It fails loudly on the two things that would sink the submission:
a draft marker that survived, and a placeholder that was never filled.
"""
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("python-pptx not installed:  python -m pip install python-pptx")

DECK = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FathomWave-Deck.pptx")

# Text that must never reach a reviewer.
BANNED = [
    "DRAFT", "TBC", "TO BE CONFIRMED", "TO BE VERIFIED", "TO BE ADDED",
    "TO BE COMPLETED", "TO BE REPLACED", "TO BE SOURCED", "TO BE INSERTED",
    "TO BE STATED", "ILLUSTRATIVE FIGURES", "PLACEHOLDER", "LOREM", "TODO", "FIXME",
]
# Organisation names that must not appear on the customers slide (requirement 4:
# "no specific names but characterizations").
# NOTE: if you reorder slides, change this or the requirement-4 check silently
# stops testing anything.
CUSTOMER_SLIDE = 7
ORG_NAMES = [
    "BIA", "Denali", "ANTHC", "NOAA", "USACE", "DOT&PF", "DCRA", "NSF",
    "AOOS", "IOOS", "EOMAP", "TCarta", "ARGANS", "Fugro", "DHI",
]

REQUIREMENTS = [
    "Description of the solution", "What ocean problem are you solving",
    "Market size", "Customers (characterizations)", "Stage of development / TRL",
    "Founder, team photo, description, advisors", "Two development + two business goals",
    "Partners", "Funded to date + next 12 months",
]


def safe(t):
    """Console-safe: Windows terminals are cp1252 and choke on the deck's typography."""
    return str(t).encode("ascii", "replace").decode("ascii")


def walk(shapes):
    """Flatten grouped shapes. Editing in PowerPoint creates groups; without this
    they hide their contents from every check below."""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for inner in walk(sh.shapes):
                yield inner
        else:
            yield sh


def iter_text(slide):
    """Yield every string on a slide, including table cells and grouped shapes."""
    for sh in walk(slide.shapes):
        if getattr(sh, "has_text_frame", False):
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    yield run, run.text
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            yield run, run.text


def main():
    if not os.path.exists(DECK):
        sys.exit("not found: " + DECK)
    prs = Presentation(DECK)
    fails, warns = [], []

    # --- format ---------------------------------------------------------
    n = len(prs.slides._sldIdLst)
    w_in, h_in = prs.slide_width / 914400.0, prs.slide_height / 914400.0
    if n > 10:
        fails.append("PAGE COUNT: %d slides, limit is 10" % n)
    if abs(w_in / h_in - 16.0 / 9.0) > 0.01:
        fails.append("ASPECT: %.2f x %.2f in is not 16:9" % (w_in, h_in))

    # --- content --------------------------------------------------------
    CANVAS_W, CANVAS_H = 960.0, 540.0   # points
    small, banned_hits, brackets, blanks, org_hits, unsized, offcanvas = [], [], [], [], [], [], []
    for i, slide in enumerate(prs.slides, 1):
        for run, text in iter_text(slide):
            if run.font.size is None:
                if text.strip():
                    unsized.append((i, text[:44]))
            elif run.font.size < Pt(12):
                small.append((i, run.font.size.pt, text[:40]))
            up = text.upper()
            for b in BANNED:
                if b in up:
                    banned_hits.append((i, b, text.strip()[:60]))
            if re.search(r"\[\s*\]|\[\s+\]|\[[^\]]{0,40}\]", text):
                brackets.append((i, text.strip()[:60]))
            if "____" in text:
                blanks.append((i, text.strip()[:60]))
            if i == CUSTOMER_SLIDE:
                for name in ORG_NAMES:
                    if re.search(r"\b" + re.escape(name) + r"\b", text):
                        org_hits.append((i, name, text.strip()[:50]))

    # Filling placeholders lengthens text; these boxes autofit by growing, so a
    # longer string can push a shape off the slide with no visible warning.
    for i, slide in enumerate(prs.slides, 1):
        for sh in walk(slide.shapes):
            if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
                continue
            l, t = sh.left / 12700.0, sh.top / 12700.0
            r, b = l + sh.width / 12700.0, t + sh.height / 12700.0
            if r > CANVAS_W + 1 or b > CANVAS_H + 1 or l < -1 or t < -1:
                offcanvas.append((i, round(r, 1), round(b, 1)))

    if small:
        fails.append("FONT FLOOR: %d run(s) below 12pt" % len(small))
    if offcanvas:
        fails.append("OFF CANVAS: %d shape(s) spill past the slide edge" % len(offcanvas))
    if unsized:
        fails.append("INHERITED FONT SIZE: %d run(s) have no explicit size (cannot prove >= 12pt)" % len(unsized))
    if banned_hits:
        fails.append("DRAFT MARKERS: %d occurrence(s) still in the deck" % len(banned_hits))
    if brackets:
        fails.append("PLACEHOLDERS: %d bracketed slot(s) unfilled" % len(brackets))
    if blanks:
        fails.append("BLANK FORM LINES: %d underscore fill line(s) remain" % len(blanks))
    if org_hits:
        fails.append("REQ 4 BREACH: %d org name(s) on the customers slide" % len(org_hits))

    # --- report ---------------------------------------------------------
    print("=" * 68)
    print("FathomWave deck pre-flight")
    print("=" * 68)
    print("slides       : %d / 10" % n)
    print("aspect       : %.2f x %.2f in  (%s)" % (w_in, h_in, "16:9 OK" if not any("ASPECT" in f for f in fails) else "WRONG"))
    print("font floor   : %s" % ("OK - nothing below 12pt" if not small else "FAIL"))
    print()

    for label, rows, fmt in (
        ("DRAFT MARKERS", banned_hits, "  slide %s  %-22s %s"),
        ("UNFILLED PLACEHOLDERS", brackets, "  slide %s  %s"),
        ("BLANK FORM LINES", blanks, "  slide %s  %s"),
        ("ORG NAMES ON CUSTOMERS SLIDE", org_hits, "  slide %s  %-10s %s"),
        ("RUNS BELOW 12PT", small, "  slide %s  %spt  %s"),
        ("SHAPES OFF CANVAS (right/bottom edge, pt)", offcanvas, "  slide %s  right=%s bottom=%s"),
        ("RUNS WITH NO EXPLICIT FONT SIZE", unsized, "  slide %s  %s"),
    ):
        if rows:
            print("%s (%d)" % (label, len(rows)))
            for r in rows[:25]:
                print(safe(fmt % r))
            if len(rows) > 25:
                print("  ... and %d more" % (len(rows) - 25))
            print()

    print("-" * 68)
    if fails:
        print("NOT READY TO SUBMIT - %d blocking issue(s):" % len(fails))
        for f in fails:
            print("  x " + f)
    else:
        print("FORMAT AND CONTENT CHECKS PASS.")
    print()
    print("Still to confirm by hand (this script cannot see these):")
    print("  - text baked inside figures/*.jpg is also >= 12pt at print size")
    print("  - every number is sourced, measured, or labelled an estimate")
    print("  - all nine required items are answered; both halves of the funding question")
    print("  - the exported PDF is what you check, not the PowerPoint")
    print("  - re-run this AFTER filling placeholders, not only before")
    print("-" * 68)
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())
